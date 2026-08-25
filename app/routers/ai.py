import io
from collections import Counter, defaultdict

from fastapi import APIRouter, Depends, HTTPException, status
from fastapi.responses import StreamingResponse
from pptx import Presentation
from sqlalchemy.orm import Session

from app.access import require_workspace_role
from app.ai_client import chat as ai_chat
from app.ai_client import complete as ai_complete
from app.database import get_db
from app.dependencies import get_current_user
from app.models.block import Block, BlockType
from app.models.stack_box import StackBox
from app.models.user import User
from app.models.workspace import Workspace, WorkspaceRole
from app.rate_limit import enforce_rate_limit
from app.schemas.ai import (
    ChatRequest,
    ChatResponse,
    DocToPptRequest,
    DraftRequest,
    DraftResponse,
    EditTextRequest,
    EditTextResponse,
    FixCodeRequest,
    FixCodeResponse,
    RepoArchitectureRequest,
    RepoArchitectureResponse,
    SummarizeRequest,
    SummarizeResponse,
)

router = APIRouter(prefix="/ai", tags=["ai"])

# Every /ai/* call is billed against our OpenAI key regardless of which user
# triggers it, so cap usage per user rather than leaving it unbounded.
_AI_RATE_LIMIT = 20
_AI_RATE_WINDOW_SECONDS = 3600

_CHAT_SYSTEM_PROMPT = (
    "당신은 StackBox 문서 작성을 돕는 어시스턴트입니다. 사용자 요청에 협조적이고 "
    "간결하게 답하세요."
)


def _enforce_ai_rate_limit(current_user: User = Depends(get_current_user)) -> None:
    enforce_rate_limit(f"ai:{current_user.id}", _AI_RATE_LIMIT, _AI_RATE_WINDOW_SECONDS)


def _get_stack_box_or_404(db: Session, stack_box_id: int) -> StackBox:
    stack_box = db.get(StackBox, stack_box_id)
    if stack_box is None:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="StackBox not found")
    return stack_box


def _get_workspace_or_404(db: Session, workspace_id: int) -> Workspace:
    workspace = db.get(Workspace, workspace_id)
    if workspace is None:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Workspace not found")
    return workspace


def _build_repo_architecture_summary(db: Session, workspace: Workspace) -> str:
    """Serialize a workspace's StackBox tree (folders/pages/canvases and their
    blocks) into plain text suitable as an AI prompt input."""
    stack_boxes = (
        db.query(StackBox)
        .filter(StackBox.workspace_id == workspace.id)
        .order_by(StackBox.parent_id, StackBox.sort_order)
        .all()
    )
    if not stack_boxes:
        return f"워크스페이스 '{workspace.name}'에는 아직 StackBox가 없습니다."

    block_counts: dict[int, Counter[BlockType]] = defaultdict(Counter)
    stack_box_ids = [stack_box.id for stack_box in stack_boxes]
    for stack_box_id, block_type in (
        db.query(Block.stack_box_id, Block.type).filter(Block.stack_box_id.in_(stack_box_ids)).all()
    ):
        block_counts[stack_box_id][block_type] += 1

    children_by_parent: dict[int | None, list[StackBox]] = defaultdict(list)
    for stack_box in stack_boxes:
        children_by_parent[stack_box.parent_id].append(stack_box)

    lines = [f"워크스페이스: {workspace.name} ({workspace.slug})"]

    def describe_blocks(stack_box_id: int) -> str:
        counts = block_counts.get(stack_box_id)
        if not counts:
            return ""
        parts = ", ".join(f"{block_type.value} 블록 {count}개" for block_type, count in counts.items())
        return f" - {parts}"

    def walk(parent_id: int | None, depth: int) -> None:
        for stack_box in children_by_parent.get(parent_id, []):
            indent = "  " * depth
            lines.append(
                f"{indent}- [{stack_box.type.value}] {stack_box.name} "
                f"(id={stack_box.id}){describe_blocks(stack_box.id)}"
            )
            walk(stack_box.id, depth + 1)

    walk(None, 0)
    return "\n".join(lines)


@router.post("/summarize", response_model=SummarizeResponse)
def summarize(
    payload: SummarizeRequest,
    current_user: User = Depends(get_current_user),
    _: None = Depends(_enforce_ai_rate_limit),
) -> SummarizeResponse:
    summary = ai_complete(
        "당신은 문서를 간결하게 요약하는 어시스턴트입니다. 핵심만 3~5문장으로 요약하세요.",
        payload.text,
    )
    return SummarizeResponse(summary=summary)


@router.post("/fix-code", response_model=FixCodeResponse)
def fix_code(
    payload: FixCodeRequest,
    current_user: User = Depends(get_current_user),
    _: None = Depends(_enforce_ai_rate_limit),
) -> FixCodeResponse:
    instructions = payload.instructions or "버그를 찾아 수정하고 개선하세요."
    language_hint = f" (언어: {payload.language})" if payload.language else ""
    system_prompt = (
        "당신은 숙련된 소프트웨어 엔지니어입니다. 아래 형식을 정확히 지켜 응답하세요:\n"
        "```\n<수정된 코드>\n```\n설명: <한두 문장 설명>"
    )
    user_prompt = f"지시사항: {instructions}{language_hint}\n\n코드:\n{payload.code}"
    raw = ai_complete(system_prompt, user_prompt)

    fixed_code = raw
    explanation = ""
    if "```" in raw:
        parts = raw.split("```")
        if len(parts) >= 2:
            code_block = parts[1]
            fixed_code = code_block.split("\n", 1)[-1] if "\n" in code_block else code_block
            remainder = "```".join(parts[2:])
            explanation = remainder.replace("설명:", "").strip()

    return FixCodeResponse(fixed_code=fixed_code.strip(), explanation=explanation)


@router.post("/edit-text", response_model=EditTextResponse)
def edit_text(
    payload: EditTextRequest,
    current_user: User = Depends(get_current_user),
    _: None = Depends(_enforce_ai_rate_limit),
) -> EditTextResponse:
    # Distinct from /fix-code on purpose: that endpoint's system prompt casts
    # the model as a software engineer and always wraps its reply in a code
    # fence, which is the wrong voice and format for revising prose selected
    # in the document editor.
    system_prompt = (
        "당신은 사용자가 선택한 텍스트를 지시사항에 따라 수정하는 어시스턴트입니다. "
        "수정된 텍스트만 출력하고, 다른 설명이나 따옴표, 코드 블록은 포함하지 마세요."
    )
    user_prompt = f"지시사항: {payload.instructions}\n\n텍스트:\n{payload.text}"
    edited_text = ai_complete(system_prompt, user_prompt).strip()
    return EditTextResponse(edited_text=edited_text)


@router.post("/draft", response_model=DraftResponse)
def draft(
    payload: DraftRequest,
    current_user: User = Depends(get_current_user),
    _: None = Depends(_enforce_ai_rate_limit),
) -> DraftResponse:
    text = ai_complete(
        "당신은 문서 초안을 작성하는 어시스턴트입니다. 마크다운 형식으로 초안을 작성하세요.",
        payload.prompt,
    )
    return DraftResponse(draft=text)


@router.post("/chat", response_model=ChatResponse)
def chat(
    payload: ChatRequest,
    current_user: User = Depends(get_current_user),
    _: None = Depends(_enforce_ai_rate_limit),
) -> ChatResponse:
    # payload.messages can only carry "user"/"assistant" roles (see ChatMessage),
    # so the system prompt below is always the first and only "system" message
    # sent to OpenAI.
    messages = [{"role": "system", "content": _CHAT_SYSTEM_PROMPT}]
    messages += [{"role": m.role, "content": m.content} for m in payload.messages]
    reply = ai_chat(messages)
    return ChatResponse(reply=reply)


@router.post("/repo-architecture", response_model=RepoArchitectureResponse)
def repo_architecture(
    payload: RepoArchitectureRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
    _: None = Depends(_enforce_ai_rate_limit),
) -> RepoArchitectureResponse:
    workspace = _get_workspace_or_404(db, payload.workspace_id)
    require_workspace_role(db, workspace.id, current_user, WorkspaceRole.viewer)

    summary = _build_repo_architecture_summary(db, workspace)

    analysis = ai_complete(
        "당신은 소프트웨어 문서/캔버스 구조를 분석하는 아키텍처 어시스턴트입니다. 아래는 StackBox "
        "워크스페이스의 구조(폴더·페이지·캔버스 계층과 각 문서에 포함된 블록 구성)를 요약한 내용입니다. "
        "이 구조를 바탕으로 전체 구성, 계층 관계, 콘텐츠 분포의 특징을 설명하는 서술형 아키텍처 분석을 "
        "작성하세요.",
        summary,
    )
    return RepoArchitectureResponse(analysis=analysis)


@router.post("/doc-to-ppt")
def doc_to_ppt(
    payload: DocToPptRequest,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
    _: None = Depends(_enforce_ai_rate_limit),
) -> StreamingResponse:
    stack_box = _get_stack_box_or_404(db, payload.stack_box_id)
    require_workspace_role(db, stack_box.workspace_id, current_user, WorkspaceRole.viewer)

    blocks = (
        db.query(Block)
        .filter(Block.stack_box_id == stack_box.id, Block.type == BlockType.markdown)
        .order_by(Block.sort_order)
        .all()
    )
    document_text = "\n\n".join(block.content for block in blocks) or stack_box.description

    outline = ai_complete(
        "당신은 문서를 PPT 슬라이드 개요로 변환하는 어시스턴트입니다. "
        "각 슬라이드는 '# ' 로 시작하는 제목 한 줄과, 그 아래 '- ' 로 시작하는 bullet point 여러 줄로 구성하세요. "
        "다른 설명 없이 개요만 출력하세요.",
        document_text,
    )

    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]

    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = stack_box.name

    current_slide = None
    for line in outline.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            current_slide = presentation.slides.add_slide(bullet_layout)
            current_slide.shapes.title.text = stripped.lstrip("#").strip()
        elif stripped.startswith("-"):
            if current_slide is None:
                current_slide = presentation.slides.add_slide(bullet_layout)
            body = current_slide.placeholders[1].text_frame
            if body.text:
                body.add_paragraph().text = stripped.lstrip("-").strip()
            else:
                body.text = stripped.lstrip("-").strip()

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)

    filename = f"{stack_box.name}.pptx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
